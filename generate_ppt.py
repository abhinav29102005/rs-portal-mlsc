from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()

# Define Thapar Maroon Color
THAPAR_MAROON = RGBColor(128, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GREY = RGBColor(64, 64, 64)

def apply_theme(slide):
    # Set background color to very light grey
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

def format_title(title_shape):
    if not title_shape: return
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = THAPAR_MAROON
        paragraph.font.bold = True
        paragraph.font.name = 'Arial'

def format_body(body_shape, size=24):
    if not body_shape: return
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = DARK_GREY
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(size)

def add_slide(title_text, body_text_list, layout_idx=1, font_size=24):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    format_title(title)
    
    if layout_idx == 1 and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear() # Clear existing text
        for idx, text in enumerate(body_text_list):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = text
        format_body(body, size=font_size)
    return slide

def add_image_slide(title_text, image_path, caption_text=""):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    format_title(title)
    
    try:
        # Add image centered
        left = Inches(1.5)
        top = Inches(2)
        height = Inches(4.5)
        slide.shapes.add_picture(image_path, left, top, height=height)
        
        if caption_text:
            txBox = slide.shapes.add_textbox(Inches(1.5), Inches(6.8), Inches(7), Inches(0.5))
            tf = txBox.text_frame
            tf.text = caption_text
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].font.color.rgb = DARK_GREY
            tf.paragraphs[0].font.italic = True
    except Exception as e:
        print(f"Failed to add image {image_path}: {e}")
        
    return slide

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_theme(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Thapar Research & Mentorship Portal"
format_title(title)
title.text_frame.paragraphs[0].font.size = Pt(44)

subtitle.text = "Bridging the Gap Between Faculty & Students\nEmpowering Innovation through Seamless Collaboration"
subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_GREY
subtitle.text_frame.paragraphs[0].font.bold = True

# Slide 2
add_slide("Introduction: The Need for Innovation", [
    "Thapar Institute thrives on cutting-edge research.",
    "However, connecting the right students with the right faculty can be challenging.",
    "We need a system that matches academic rigor with technological efficiency.",
    "The Research Portal is designed to be that bridge."
])

# Slide 3
add_slide("The Current Problem: Fragmented Systems", [
    "Faculty receive hundreds of unorganized emails.",
    "Students struggle to find relevant active projects.",
    "No centralized repository of ongoing research.",
    "Time wasted on logistics instead of actual research."
])

# Slide 4
add_slide("The Proposed Solution: A Centralized Hub", [
    "A single unified platform for all research activities.",
    "Streamlined discovery of projects and mentorships.",
    "Automated tracking of applications and statuses.",
    "Designed specifically for Thapar's ecosystem."
])

# Slide 5
add_slide("Key Objectives of the Portal", [
    "Increase student participation in high-impact research.",
    "Reduce administrative overhead for faculty and PIs.",
    "Provide transparent metrics on research engagement.",
    "Foster a culture of interdisciplinary collaboration."
])

# Slide 6
add_slide("Novelty 1: Unified Ecosystem", [
    "No more scattered Google Forms or Excel sheets.",
    "Everything from project ideation to final selection happens in one place.",
    "A clean, modern interface that users actually want to use."
])

# Slide 7
add_slide("Novelty 2: Automated Mentorship Alignment", [
    "The system doesn't just list projects; it aligns them.",
    "Faculty can set strict prerequisite skills for projects.",
    "Only eligible students can apply, ensuring high-quality candidates."
])

# Slide 8
add_slide("Novelty 3: Data-Driven Matching", [
    "Leveraging student profiles (CGPA, Tech Stack, Previous Projects).",
    "Providing faculty with instant analytics on an applicant's suitability.",
    "Smart tagging system for rapid sorting and filtering."
])

# Slide 9
add_slide("Methodology: Phase 1 - Data Aggregation", [
    "Consolidate existing faculty profiles and research domains.",
    "Migrate legacy project data into a modern, structured format.",
    "Ensure data integrity and consistency across departments."
])

# Slide 10
add_slide("Methodology: Phase 2 - Secure Access", [
    "Implement robust authentication for both students and faculty.",
    "Role-based access control (RBAC) to protect sensitive data.",
    "Seamless integration with institutional credentials."
])

# Slide 11
add_slide("Methodology: Phase 3 - Application Workflow", [
    "1. Faculty posts a project with required skills.",
    "2. Students browse, filter, and apply.",
    "3. Faculty reviews applications via a specialized dashboard.",
    "4. One-click accept/reject with automated notifications."
])

# Slide 12: Image 1
add_image_slide("Portal Solution: Faculty Dashboard", 
                "/home/bigboyaks/.gemini/antigravity/brain/0dd69ef7-a9d2-4b61-a03b-1d6bf1fe0dc0/portal_dashboard_1783447262284.png", 
                "A clean, centralized view of ongoing projects, pending applications, and recent activity.")

# Slide 13: Image 2
add_image_slide("Portal Solution: Project Discovery", 
                "/home/bigboyaks/.gemini/antigravity/brain/0dd69ef7-a9d2-4b61-a03b-1d6bf1fe0dc0/student_view_1783447294664.png", 
                "Students can easily browse, filter, and apply to research projects matching their interests.")

# Slide 14: Image 3
add_image_slide("Portal Solution: Mentor Matching", 
                "/home/bigboyaks/.gemini/antigravity/brain/0dd69ef7-a9d2-4b61-a03b-1d6bf1fe0dc0/mentor_matching_1783447278537.png", 
                "Faculty can review student profiles, complete with skill tags and GPAs, for rapid decision-making.")

# Slide 15
add_slide("Tech Stack: Overview", [
    "Built for speed, scale, and maintainability.",
    "Frontend: Next.js (React) & Tailwind CSS.",
    "Backend: Cloudflare Workers (Edge Computing).",
    "Database: Cloudflare D1 (Serverless SQL).",
    "Language: TypeScript."
])

# Slide 16
add_slide("Tech Stack: Deep Dive", [
    "Next.js provides server-side rendering for instant load times.",
    "Tailwind CSS ensures a responsive, Thapar-branded aesthetic.",
    "Cloudflare D1 provides a globally distributed, low-latency database.",
    "TypeScript prevents runtime errors, ensuring a stable platform."
])

# Slide 17
add_slide("Easing Faculty Stress", [
    "No More Inbox Clutter: Move away from hundreds of emails.",
    "Pre-Screened Candidates: Only qualified students get through.",
    "Automated Responses: Bulk accept/reject with one click.",
    "Focus on Research, Not Admin: Let the platform handle logistics."
])

# Slide 18
add_slide("Future Steps: Advanced Student Filter", [
    "Allow faculty to set complex boolean filters (e.g., Python AND Machine Learning).",
    "Automated ranking of students based on project relevance.",
    "Integration with past academic performance metrics."
])

# Slide 19
add_slide("Future Steps: Dedicated Mentor Service", [
    "Expand beyond project-based matching.",
    "Allow students to request 1-on-1 career or academic mentorship.",
    "Bring alumni onboard as external mentors.",
    "Algorithmic matchmaking based on long-term career goals."
])

# Slide 20: Conclusion
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_theme(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Thank You"
format_title(title)
title.text_frame.paragraphs[0].font.size = Pt(44)

subtitle.text = "Empowering Thapar's Research Community\nReady for Questions"
subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_GREY
subtitle.text_frame.paragraphs[0].font.bold = True

prs.save('Thapar_Research_Portal_Presentation.pptx')
print('Presentation generated successfully with 20 slides: Thapar_Research_Portal_Presentation.pptx')
