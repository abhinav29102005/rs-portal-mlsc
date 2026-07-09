from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_CONNECTOR

# Initialize presentation
prs = Presentation()

# Define Thapar Maroon Color
THAPAR_MAROON = RGBColor(128, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GREY = RGBColor(64, 64, 64)

def apply_theme(slide):
    # Set background color to very light grey
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)

def format_title(title_shape):
    if not title_shape: return
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = THAPAR_MAROON
        paragraph.font.bold = True
        paragraph.font.name = 'Arial'

def format_body(body_shape, size=20):
    if not body_shape: return
    for paragraph in body_shape.text_frame.paragraphs:
        paragraph.font.color.rgb = DARK_GREY
        paragraph.font.name = 'Arial'
        paragraph.font.size = Pt(size)
        paragraph.space_after = Pt(10)

def add_slide(title_text, body_text_list, layout_idx=1, font_size=20):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    format_title(title)
    
    if layout_idx == 1 and hasattr(slide, 'placeholders') and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for idx, text in enumerate(body_text_list):
            p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            p.text = text
        format_body(body, size=font_size)
    return slide

def add_image_slide(title_text, image_path, caption_text="", layout_idx=5):
    slide_layout = prs.slide_layouts[layout_idx] 
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    format_title(title)
    
    try:
        # Add image centered
        left = Inches(1.0)
        top = Inches(1.8)
        width = Inches(8)
        slide.shapes.add_picture(image_path, left, top, width=width)
        
        if caption_text:
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = caption_text
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GREY
            p.font.italic = True
    except Exception as e:
        print(f"Failed to add image {image_path}: {e}")
        
    return slide

def add_diagram_slide(title_text):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    apply_theme(slide)
    
    title = slide.shapes.title
    title.text = title_text
    format_title(title)
    
    # Create an architectural diagram using shapes
    shapes = slide.shapes
    
    # 1. Frontend
    fe_shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(3.5), Inches(2), Inches(1))
    fe_shape.text = "Next.js Frontend\n(Students & Faculty)"
    fe_shape.fill.solid()
    fe_shape.fill.fore_color.rgb = THAPAR_MAROON
    
    # 2. API / Edge
    api_shape = shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(4), Inches(3), Inches(2), Inches(2))
    api_shape.text = "Next.js\nAPI Routes"
    api_shape.fill.solid()
    api_shape.fill.fore_color.rgb = RGBColor(100, 100, 100)
    
    # 3. Database
    db_shape = shapes.add_shape(MSO_SHAPE.CAN, Inches(7), Inches(3.5), Inches(2), Inches(1))
    db_shape.text = "Turso DB\n(LibSQL Database)"
    db_shape.fill.solid()
    db_shape.fill.fore_color.rgb = THAPAR_MAROON

    # Connectors
    connector1 = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3), Inches(4), Inches(4), Inches(4))
    connector1.line.color.rgb = DARK_GREY
    connector1.line.width = Pt(2)
    
    connector2 = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6), Inches(4), Inches(7), Inches(4))
    connector2.line.color.rgb = DARK_GREY
    connector2.line.width = Pt(2)

    # Caption
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(6), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "System Architecture: Secure, Global Edge Delivery with Serverless SQL."
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.color.rgb = DARK_GREY
    
    return slide

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_theme(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "RAMP: Research And Mentorship Portal"
format_title(title)
title.text_frame.paragraphs[0].font.size = Pt(44)

subtitle.text = "A College Project\nBridging the Gap Between Faculty & Students\nEmpowering Innovation through Seamless Collaboration"
subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_GREY
subtitle.text_frame.paragraphs[0].font.bold = True
subtitle.text_frame.paragraphs[0].font.size = Pt(24)

try:
    slide.shapes.add_picture("public/thapar-logo.png", Inches(4), Inches(0.5), width=Inches(2))
except:
    pass

# Slide 2: Index / Agenda
add_slide("Agenda & Index", [
    "1. Introduction: The Need for Innovation",
    "2. The Current Problem: Fragmented Systems",
    "3. The Proposed Solution: A Centralized Hub",
    "4. Key Objectives & Target Outcomes",
    "5. Novelty: Unified Ecosystem & Automated Alignment",
    "6. Methodology: Phased Implementation Strategy",
    "7. Real Portal Showcases (Dashboard, Discovery, Matching)",
    "8. Tech Stack Architecture (Diagram & Details)",
    "9. Value Proposition: Easing Faculty Stress",
    "10. Future Horizons: Advanced Filtering & Mentor Services",
    "11. Conclusion & Open Q&A"
], font_size=22)

# Slide 3-11: Content (same as previous script)
add_slide("Introduction: The Need for Innovation", [
    "Context: Thapar Institute of Engineering & Technology has a vibrant research culture.",
    "Challenge: Despite high enthusiasm, connecting the right students with the right faculty can be daunting.",
    "Goal: We require a system that matches academic rigor with technological efficiency to streamline workflows.",
    "Vision: RAMP (Research And Mentorship Portal) is designed to act as a definitive bridge, turning potential collaboration into tangible outputs."
])

add_slide("The Current Problem: Fragmented Systems", [
    "Information Silos: Faculty are forced to rely on hundreds of unorganized, scattered emails for applications.",
    "Visibility Issues: Students struggle to discover relevant, active projects because they aren't listed centrally.",
    "Inefficient Processes: Time is heavily wasted on logistics—coordinating meetings, tracking status, filtering resumes.",
    "Lack of Analytics: No centralized repository exists to track ongoing research metrics and student involvement across departments."
])

add_slide("The Proposed Solution: A Centralized Hub", [
    "Unified Ecosystem: A single, integrated web platform covering all research activities at the institute.",
    "End-to-End Tracking: From initial project ideation to final student selection, every step is automated and tracked.",
    "Customized Workflows: Designed specifically for Thapar's unique academic environment, respecting existing hierarchies and procedures.",
    "Data Accessibility: Provides intuitive dashboards tailored separately for faculty, students, and administrators."
])

add_slide("Key Objectives of the Portal", [
    "Maximize Engagement: Increase overall student participation in high-impact, meaningful research.",
    "Minimize Administrative Overhead: Drastically reduce the time faculty and PIs spend on managing applications.",
    "Enhance Transparency: Provide clear, real-time metrics on research engagement and application statuses.",
    "Promote Interdisciplinary Work: Foster a campus culture where cross-departmental collaboration is visible and accessible."
])

add_slide("Novelty: Unified & Automated Ecosystem", [
    "Elimination of Legacy Tools: Replaces scattered Google Forms, WhatsApp messages, and Excel sheets with a cohesive UI.",
    "Automated Mentorship Alignment: The system actively prevents 'spray-and-pray' applications.",
    "Prerequisite Enforcement: Faculty can set strict, verifiable skills for projects (e.g., must know React and Next.js).",
    "Quality Assurance: Only students who meet the rigorous criteria are allowed to submit a formal application."
])

add_slide("Novelty: Data-Driven Matching", [
    "Comprehensive Profiles: Leverages detailed student profiles including CGPA, specific Tech Stacks, and Portfolio links.",
    "Instant Analytics: Provides faculty with immediate, visual analytics indicating an applicant's exact suitability score.",
    "Smart Tagging: Utilizes a robust tagging architecture for rapid sorting, categorization, and retrieval of research domains."
])

add_slide("Methodology: Phase 1 - Data Aggregation", [
    "Consolidation: Gather and consolidate existing faculty profiles, research papers, and domains into a unified database.",
    "Data Migration: Safely migrate any legacy project data into our modern, structured, schema-validated formats.",
    "Normalization: Ensure strict data integrity, consistency, and typing across all departments using TypeScript and Zod.",
    "Indexing: Prepare the database for high-speed searches and complex queries for the frontend."
])

add_slide("Methodology: Phase 2 - Secure Access & Profiles", [
    "Robust Authentication: Implement enterprise-grade, secure authentication flows for both students and faculty.",
    "RBAC Implementation: Strict Role-Based Access Control to ensure sensitive data (like proposals) is protected.",
    "Profile Customization: Enable dynamic user profiles where users can showcase their achievements securely.",
    "Edge Delivery: Utilize Cloudflare's edge network to serve authenticated sessions with ultra-low latency."
])

add_slide("Methodology: Phase 3 - Application Workflow", [
    "1. Project Posting: Faculty post a detailed project description along with mandatory skill requirements.",
    "2. Discovery & Filtering: Students browse the central repository, filter by their tech stack, and apply.",
    "3. Centralized Review: Faculty receive applications directly in their specialized dashboard, circumventing email.",
    "4. Resolution: Faculty perform a one-click accept/reject action, triggering automated, polite notification emails."
])

# Slides 12-14: Real Portal Access using real captured screenshots
add_image_slide("Portal Showcase: Faculty Dashboard", 
                "faculty_dashboard.png", 
                "Real Portal Screenshot: A clean, centralized view demonstrating ongoing projects and pending applications for faculty.")

add_image_slide("Portal Showcase: Project Discovery (Student View)", 
                "student_discover.png", 
                "Real Portal Screenshot: The student interface allowing intuitive browsing and instant applications to relevant opportunities.")

add_image_slide("Portal Showcase: Mentor Matching", 
                "faculty_mentors.png", 
                "Real Portal Screenshot: The faculty interface for reviewing applicants, displaying profiles, and making matching decisions.")

# Slide 15: Architecture Diagram
add_diagram_slide("Tech Stack Architecture: System Flow Diagram")

# Slide 16: Tech Stack
add_slide("Tech Stack Architecture: Core Technologies", [
    "Philosophy: Built entirely for blistering speed, massive scale, and long-term maintainability.",
    "Frontend Framework: Next.js (React) utilizing the App Router for optimal Server-Side Rendering (SSR).",
    "Styling Engine: Tailwind CSS, providing a responsive, cohesive, and Thapar-branded visual aesthetic.",
    "Backend & API: Next.js API Routes, providing robust backend integration natively.",
    "Database: Turso DB (LibSQL), offering an edge-ready, ultra-fast SQLite distributed database.",
    "Language: Strict TypeScript, ensuring end-to-end type safety and drastically fewer runtime errors."
])

# Slide 17: Easing Faculty Stress
add_slide("Easing Faculty Stress: Automating the Mundane", [
    "No More Inbox Clutter: Fully transition away from hundreds of disorganized emails to a structured dashboard.",
    "Pre-Screened Quality: The system physically blocks unqualified candidates from applying, saving immense review time.",
    "Automated Responses: With bulk accept/reject capabilities, faculty no longer need to draft individual rejection emails.",
    "Focus on Core Work: By letting the platform handle the logistical heavy lifting, faculty can dedicate their time to actual research."
])

# Slide 18: Future Steps 1
add_slide("Future Steps: Advanced Student Filtering", [
    "Complex Logic: Allow faculty to set complex boolean filters (e.g., must know 'Python' AND 'Machine Learning', but NOT 'Java').",
    "Automated Ranking: The portal will automatically rank incoming students based on how closely their profiles match the project.",
    "Academic Integration: Direct integration with institutional records to automatically verify CGPAs and course completions.",
    "Portfolio Parsing: AI-driven analysis of a student's GitHub or past projects to determine authentic competency."
])

# Slide 19: Future Steps 2
add_slide("Future Steps: Dedicated Mentor Service", [
    "Beyond Projects: Expand the portal's scope beyond strict project-based matching to general mentorship.",
    "1-on-1 Guidance: Allow students to formally request 1-on-1 career, academic, or domain-specific mentorship from faculty.",
    "Alumni Network: Bring esteemed alumni onboard as external mentors to guide current students entering the industry.",
    "Algorithmic Pairing: Implement machine-learning matchmaking to pair students with mentors based on long-term career goals."
])

# Slide 20: Conclusion
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
apply_theme(slide)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Thank You"
format_title(title)
title.text_frame.paragraphs[0].font.size = Pt(44)

subtitle.text = "Empowering Thapar's Research Community\nReady for Questions and Feedback"
subtitle.text_frame.paragraphs[0].font.color.rgb = DARK_GREY
subtitle.text_frame.paragraphs[0].font.bold = True
subtitle.text_frame.paragraphs[0].font.size = Pt(24)

try:
    slide.shapes.add_picture("public/thapar-logo.png", Inches(4), Inches(0.5), width=Inches(2))
except:
    pass

prs.save('Thapar_Research_Portal_Presentation_Detailed.pptx')
print('Detailed presentation generated successfully with 20 slides, real screenshots, and diagrams.')
