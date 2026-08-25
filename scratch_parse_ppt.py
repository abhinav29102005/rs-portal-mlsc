import sys

def parse_pptx(filepath):
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            print(f"--- SLIDE {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except ImportError:
        print("python-pptx is not installed. Please run 'pip install python-pptx'")
    except Exception as e:
        print(f"Error reading file: {e}")

if __name__ == "__main__":
    if len(sys.argv) > 1:
        parse_pptx(sys.argv[1])
    else:
        print("Please provide a filepath")
